"""
@FileName：pptx_office.py\n
@Description：\n
@Author：道长\n
@Time：2021/3/15 13:39\n
@Department：运营部\n
@Website：www.geekaso.com.com\n
@Copyright：©2019-2021 七麦数据
"""
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.shapes.group import GroupShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from qm_spider import *


# 设置自动检测，检测文字并操作；
class Text_Frame_Detect:
    def __init__(self, shapes):
        self.shapes = shapes

    def shape_textOne_detect(self, match_old_text='', match_new_text=''):
        """
            * 使用传入的当前页，检测当前页中的单个文本框；
        """
        for shape in self.shapes:
            if shape.has_text_frame:
                if shape.has_text_frame:
                    for x in shape.text_frame.paragraphs:
                        for i in x.runs:
                            if i.text == match_old_text:
                                i.text = match_new_text

    def shape_textMore_detect(self, match_old_text='', match_new_text=''):
        """
            * 使用传入的当前页，检测当前页中的嵌套文本框；
        """
        for shape in self.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shs in shape.shapes:
                    if shs.has_text_frame:
                        for x in shs.text_frame.paragraphs:
                            for i in x.runs:
                                if i.text == match_old_text:
                                    i.text = match_new_text

    def shape_id_detect(self):
        """
            * 使用传入的当前页，检测当前页中的内容ID；
        """
        pass

# 设置自动检测，检测图像并操作；
class Image_Frame_Detect:
    def __init__(self, shapes):
        self.shapes = shapes

    def image_detect(self):
        for shape in self.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(shape.image.sha1)

